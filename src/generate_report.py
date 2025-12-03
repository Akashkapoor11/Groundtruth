import os, json
from jinja2 import Template
from pathlib import Path

REPORT_JSON = os.path.join(os.path.dirname(__file__), "..", "data", "final", "insights_ai.json")
HTML_TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "assets", "template.html")
OUT_PDF = os.path.join(os.path.dirname(__file__), "..", "data", "final", "report.pdf")
OUT_PPTX = os.path.join(os.path.dirname(__file__), "..", "data", "final", "report.pptx")

def run():
    with open(REPORT_JSON) as f:
        payload = json.load(f)
    analysis = payload.get('analysis', {})
    insights = payload.get('insights', {})

    # Table rows for HTML
    table_rows = []
    for m in analysis.get('metrics', []):
        table_rows.append({
            'week': m.get('week'),
            'visitors': m.get('visitors'),
            'sales': m.get('sales'),
            'revenue': m.get('revenue'),
        })

    anomalies = analysis.get("anomalies", [])

    tpl_text = Path(HTML_TEMPLATE).read_text()
    tpl = Template(tpl_text)
    summary_text = (
        insights.get("ai_text")
        if isinstance(insights, dict) and "ai_text" in insights
        else insights.get("summary", "")
    )

    html = tpl.render(
        title="Automated Insight Report",
        period="All weeks",
        summary=summary_text,
        table_rows=table_rows,
        anomalies=anomalies,
        recommendations="\n".join(insights.get("recommendations", [])),
    )

    # PDF Export (will fail on Windows — fallback is HTML)
    try:
        from weasyprint import HTML
        HTML(string=html).write_pdf(OUT_PDF)
        print("Generated PDF report at", OUT_PDF)
    except Exception as e:
        print("WeasyPrint failed:", e)
        html_out = os.path.join(os.path.dirname(__file__), "..", "data", "final", "report.html")
        Path(html_out).write_text(html)
        print("Saved HTML report to", html_out)

    # PPTX Export (fixed version)
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Slide 1 — Title + summary
        title_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_slide.shapes.title.text = "Automated Insight Report"
        title_slide.placeholders[1].text = summary_text

        # Slide 2 — Table
        table_slide = prs.slides.add_slide(prs.slide_layouts[5])

        rows_count = len(table_rows) + 1  # header + data
        cols = 4

        table_shape = table_slide.shapes.add_table(
            rows_count, cols, Inches(0.5), Inches(0.5), Inches(9.0), Inches(5.0)
        )
        table = table_shape.table

        headers = ["Week", "Visitors", "Sales", "Revenue"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h

        for i, row in enumerate(table_rows, start=1):
            table.cell(i, 0).text = str(row["week"])
            table.cell(i, 1).text = str(row["visitors"])
            table.cell(i, 2).text = str(row["sales"])
            table.cell(i, 3).text = str(row["revenue"])

        prs.save(OUT_PPTX)
        print("Generated PPTX report at", OUT_PPTX)

    except Exception as e:
        print("PPTX generation failed:", e)

    return True

if __name__ == "__main__":
    run()
